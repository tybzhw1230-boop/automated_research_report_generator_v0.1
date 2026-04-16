from __future__ import annotations

import math
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Type

from crewai.tools import BaseTool
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


FALLBACK_TEXT = "缺乏信息"
PRIMARY_FONT = "Microsoft YaHei"
TITLE_FONT = "SimHei"

SLIDE_WIDTH = Inches(11.69)
SLIDE_HEIGHT = Inches(8.27)

PAGE_LEFT = Inches(0.40)
PAGE_TOP = Inches(0.20)
PAGE_WIDTH = Inches(10.89)

TITLE_TOP = Inches(0.18)
TITLE_HEIGHT = Inches(0.38)
TITLE_RULE_TOP = Inches(0.68)
BANNER_TOP = Inches(0.82)
BANNER_HEIGHT = Inches(0.34)

TOP_PANEL_TOP = Inches(1.28)
TOP_PANEL_HEIGHT = Inches(2.58)
PANEL_GAP = Inches(0.22)
TOP_PANEL_WIDTH = Inches(5.335)

HIGHLIGHT_TOP = Inches(4.12)
HIGHLIGHT_HEIGHT = Inches(1.90)
RISK_TOP = Inches(6.20)
RISK_HEIGHT = Inches(1.30)

BACKGROUND = RGBColor(255, 255, 255)
TITLE_TEXT = RGBColor(32, 27, 23)
PRIMARY_RED = RGBColor(151, 49, 29)
SECONDARY_GOLD = RGBColor(198, 149, 82)
LABEL_BG = RGBColor(240, 233, 224)
PANEL_BG = RGBColor(252, 249, 246)
ROW_ODD_BG = RGBColor(247, 242, 236)
ROW_EVEN_BG = RGBColor(241, 235, 229)
PANEL_BORDER = RGBColor(208, 177, 130)
MUTED_TEXT = RGBColor(86, 77, 70)

REQUIRED_FINANCIAL_LABELS = (
    "营业收入",
    "营业收入增长率",
    "毛利率",
    "净利率",
    "净利润增长率",
    "总资产",
    "资产负债率",
    "ROE",
)

DIRECT_FINANCIAL_LABEL_ALIASES = {
    "营业收入": ("营业收入", "收入", "营业总收入"),
    "营业收入增长率": (
        "营业收入增长率",
        "营业收入同比增长率",
        "收入增长率",
        "收入同比增长率",
        "营收增长率",
        "营收同比增长率",
    ),
    "毛利率": ("毛利率", "销售毛利率", "毛利率(%)"),
    "净利率": ("净利率", "归母净利率", "销售净利率", "净利率(%)"),
    "净利润增长率": (
        "净利润增长率",
        "归母净利润增长率",
        "净利润同比增长率",
        "归母净利润同比增长率",
    ),
    "总资产": ("总资产", "资产总计"),
    "资产负债率": ("资产负债率", "负债率"),
    "ROE": ("ROE", "净资产收益率", "加权平均净资产收益率"),
}

SUPPORTING_FINANCIAL_LABEL_ALIASES = {
    "净利润": ("净利润", "归母净利润", "本公司拥有人应占利润", "归母利润"),
    "股东权益": ("股东权益", "所有者权益", "归母净资产", "归属于母公司股东权益"),
}


class InvestmentSnapshotFinancialRow(BaseModel):
    """
    目的：定义单页快照里财务表的标准行结构。
    功能：统一承载行名与按展示期间对齐后的值列表。
    实现逻辑：先在工具内部解析或推导完整行，再把结果收口到这个模型交给渲染层。
    可调参数：`label` 为展示行名，`values` 为按期间顺序排列的值。
    默认参数及原因：字段全部必填，原因是表格渲染不应在最后阶段再猜测结构。
    """

    label: str = Field(..., description="财务指标名称")
    values: list[str] = Field(..., description="与展示期间一一对应的值列表")


class InvestmentSnapshotOverviewProductItem(BaseModel):
    """
    目的：定义公司概况区的产品条目结构。
    功能：把产品名和说明拆开，便于工具在 PPT 中分别设置层级和样式。
    实现逻辑：Agent 只需按 `name + description` 输出，具体排版由工具统一完成。
    可调参数：`name` 为产品名，`description` 为一句话说明。
    默认参数及原因：字段全部必填，原因是概况区固定需要完整的“名称 + 描述”对。
    """

    name: str = Field(..., description="产品名称")
    description: str = Field(..., description="产品一句话说明")


class InvestmentSnapshotTitledItem(BaseModel):
    """
    目的：定义亮点和风险条目的统一结构。
    功能：把短标题和详细说明拆开，便于工具用粗细和颜色做层级区分。
    实现逻辑：Agent 输出结构化字段，工具按统一模板排成短段落。
    可调参数：`title` 为条目标题，`detail` 为条目说明。
    默认参数及原因：字段全部必填，原因是亮点和风险都需要可扫描的标题层。
    """

    title: str = Field(..., description="条目标题")
    detail: str = Field(..., description="条目详情")


class InvestmentSnapshotPptInput(BaseModel):
    """
    目的：定义 investment snapshot 工具的稳定入参。
    功能：把 PPT 导出需要的叙事内容、输出路径和财务原文统一约束为结构化字段。
    实现逻辑：Agent 负责压缩内容，工具负责解析财务表并实际绘制 PPT。
    可调参数：输出路径、标题、摘要、产品项、财务原文、亮点和风险。
    默认参数及原因：所有核心字段都设为必填，原因是该工具要独立完成最终导出。
    """

    pptx_path: str = Field(..., description="输出 PPTX 文件路径")
    slide_title: str = Field(..., description="页面主标题")
    positioning_line: str = Field(..., description="顶部红色横幅中的定位短句")
    overview_summary: str = Field(..., description="公司概况摘要")
    overview_product_items: list[InvestmentSnapshotOverviewProductItem] = Field(
        ...,
        description="公司概况中的产品条目列表",
    )
    financial_source_markdown: str = Field(
        ...,
        description="财务专题 pack 的完整 Markdown 原文",
    )
    highlight_items: list[InvestmentSnapshotTitledItem] = Field(
        ...,
        description="投资亮点条目列表",
    )
    risk_items: list[InvestmentSnapshotTitledItem] = Field(
        ...,
        description="投资风险条目列表",
    )


@dataclass(slots=True)
class ParsedFinancialTable:
    """
    目的：承载从 Markdown 财务表中解析出的关键结果。
    功能：保存全量期间、展示期间、展示列索引、单位说明和行映射。
    实现逻辑：解析层一次性输出标准结果，后续推导和渲染都复用这份结构。
    可调参数：各字段由解析函数写入，不直接对外暴露额外开关。
    默认参数及原因：不提供运行时默认值，原因是来源信息必须来自真实解析结果。
    """

    periods: list[str]
    display_periods: list[str]
    selected_indexes: list[int]
    unit_note: str
    row_lookup: dict[str, list[str]]


class InvestmentSnapshotPptTool(BaseTool):
    """
    目的：把上游分析 pack 压缩为一页横向 A4 的投资快照 PPT。
    功能：负责财务表解析、8 个固定指标映射、版式绘制和 PPTX 文件导出。
    实现逻辑：先规范化 Agent 传入内容，再从 finance pack 中解析最近三期，最后绘制单页并保存。
    可调参数：全部来自 `InvestmentSnapshotPptInput` 定义的字段。
    默认参数及原因：工具名和 schema 固定，原因是写作任务已围绕这套接口写死提示词。
    """

    name: str = "investment_snapshot_ppt_tool"
    description: str = (
        "Create a single-slide A4 landscape investment snapshot PowerPoint using "
        "structured summary content and finance pack markdown."
    )
    args_schema: Type[BaseModel] = InvestmentSnapshotPptInput

    _PERIOD_PATTERN = re.compile(
        r"^(ttm|fy[- ]?\d+|fq[- ]?\d+|fq0/fy0|20\d{2}(a|e)?|20\d{2}q[1-4](a|e)?|20\d{2}h[12](a|e)?)$",
        flags=re.IGNORECASE,
    )

    def _run(
        self,
        pptx_path: str,
        slide_title: str,
        positioning_line: str,
        overview_summary: str,
        overview_product_items: list[InvestmentSnapshotOverviewProductItem],
        financial_source_markdown: str,
        highlight_items: list[InvestmentSnapshotTitledItem],
        risk_items: list[InvestmentSnapshotTitledItem],
    ) -> str:
        """
        目的：执行 investment snapshot PPT 的完整生成流程。
        功能：校正输入内容、解析财务表、绘制一页 PPT 并保存到目标路径。
        实现逻辑：先创建输出目录，再调用解析与渲染方法，最后返回可直接写入任务结果的成功消息。
        可调参数：全部来自任务层传入的输出路径、摘要、结构化条目和财务原文。
        默认参数及原因：无额外隐式默认行为，原因是要避免悄悄改写 Agent 已确认的输入。
        """

        output_file = Path(pptx_path).expanduser().resolve()
        output_file.parent.mkdir(parents=True, exist_ok=True)

        try:
            parsed_table = self._parse_financial_markdown(financial_source_markdown)
        except Exception as exc:
            raise ValueError(
                "Failed to build investment snapshot PPT during finance parsing "
                f"for output path {output_file}: {exc}"
            ) from exc

        presentation = self._build_presentation(
            company_name=self._derive_company_name_from_output_path(output_file),
            slide_title=self._normalize_text(slide_title, fallback_text="投资要点速览"),
            positioning_line=self._normalize_text(positioning_line, fallback_text=FALLBACK_TEXT),
            overview_summary=self._normalize_text(overview_summary, fallback_text=FALLBACK_TEXT),
            overview_product_items=self._normalize_overview_product_items(overview_product_items),
            financial_table=parsed_table,
            financial_rows=self._build_snapshot_financial_rows(parsed_table),
            highlight_items=self._normalize_titled_items(highlight_items, min_items=3, max_items=5, fallback_title="投资亮点"),
            risk_items=self._normalize_titled_items(risk_items, min_items=1, max_items=2, fallback_title="投资风险"),
        )
        presentation.save(str(output_file))
        return f"PPT created successfully at: {output_file}"

    def _derive_company_name_from_output_path(self, pptx_path: Path) -> str:
        """
        目的：从约定好的输出文件名推导公司名称。
        功能：避免在当前任务链里额外增加 `company_name` 工具入参。
        实现逻辑：优先剥离 `_investment_snapshot` 后缀，剥离失败时退回文件 stem。
        可调参数：`pptx_path` 为当前导出目标路径。
        默认参数及原因：未命中约定命名时退回 stem，原因是这样最稳妥且不引入额外依赖。
        """

        stem = pptx_path.stem.strip()
        suffix = "_investment_snapshot"
        if stem.endswith(suffix):
            trimmed = stem[: -len(suffix)].strip()
            if trimmed:
                return trimmed
        return stem or "未知公司"

    def _normalize_text(self, text: str | None, fallback_text: str) -> str:
        """
        目的：统一清洗 Agent 传入的短文本。
        功能：压缩多余空白，并在空值时回落到显式兜底文案。
        实现逻辑：只做最小清洗，不裁剪长度也不改写语义。
        可调参数：`text` 为原始文本，`fallback_text` 为缺失时的兜底值。
        默认参数及原因：兜底文案由调用方显式传入，原因是不同区域的占位语不同。
        """

        normalized_text = " ".join(str(text or "").split()).strip()
        return normalized_text or fallback_text

    def _normalize_overview_product_items(
        self,
        items: list[InvestmentSnapshotOverviewProductItem] | list[dict[str, object]],
    ) -> list[InvestmentSnapshotOverviewProductItem]:
        """
        目的：把概况区的产品项规范成固定 3 条。
        功能：兼容 CrewAI 可能传入的 dict 或模型对象，并补齐缺失条目。
        实现逻辑：先提取有效项，再不足 3 条时补固定占位项，超过 3 条时截断。
        可调参数：`items` 为 Agent 传入的产品条目列表。
        默认参数及原因：固定 3 条，原因是当前 A4 单页布局就是按 3 条做的空间预算。
        """

        normalized_items: list[InvestmentSnapshotOverviewProductItem] = []
        for raw_item in items[:3]:
            name, description = self._extract_named_item_payload(raw_item)
            if not name and not description:
                continue
            normalized_items.append(
                InvestmentSnapshotOverviewProductItem(
                    name=self._normalize_text(name, fallback_text="核心产品"),
                    description=self._normalize_text(description, fallback_text=FALLBACK_TEXT),
                )
            )
        while len(normalized_items) < 3:
            normalized_items.append(
                InvestmentSnapshotOverviewProductItem(
                    name="核心产品",
                    description=FALLBACK_TEXT,
                )
            )
        return normalized_items[:3]

    def _normalize_titled_items(
        self,
        items: list[InvestmentSnapshotTitledItem] | list[dict[str, object]],
        *,
        min_items: int,
        max_items: int,
        fallback_title: str,
    ) -> list[InvestmentSnapshotTitledItem]:
        """
        目的：统一规范亮点和风险条目列表。
        功能：兼容 dict / 模型对象输入，并按布局要求限制条目数。
        实现逻辑：先提取 title/detail，再截断、补位并保持每条只有一个标题。
        可调参数：条目列表、最小条数、最大条数和缺失时的占位标题。
        默认参数及原因：亮点和风险的条数约束来自当前任务提示词，因此由调用方显式指定。
        """

        normalized_items: list[InvestmentSnapshotTitledItem] = []
        for raw_item in items[:max_items]:
            title, detail = self._extract_titled_item_payload(raw_item)
            if not title and not detail:
                continue
            normalized_items.append(
                InvestmentSnapshotTitledItem(
                    title=self._normalize_text(title, fallback_text=fallback_title),
                    detail=self._normalize_text(detail, fallback_text=FALLBACK_TEXT),
                )
            )
        while len(normalized_items) < min_items:
            normalized_items.append(
                InvestmentSnapshotTitledItem(
                    title=fallback_title,
                    detail=FALLBACK_TEXT,
                )
            )
        return normalized_items[:max_items]

    def _extract_named_item_payload(
        self,
        item: InvestmentSnapshotOverviewProductItem | dict[str, object],
    ) -> tuple[str, str]:
        """
        目的：兼容不同来源的产品条目对象。
        功能：统一抽取 `name` 和 `description` 两个字段。
        实现逻辑：优先读取模型对象属性，若是字典则按键名读取。
        可调参数：`item` 为单个产品项。
        默认参数及原因：无法识别时返回空字符串，原因是让上层统一做补位更稳。
        """

        if isinstance(item, InvestmentSnapshotOverviewProductItem):
            return str(item.name).strip(), str(item.description).strip()
        if isinstance(item, dict):
            return str(item.get("name", "")).strip(), str(item.get("description", "")).strip()
        return "", ""

    def _extract_titled_item_payload(
        self,
        item: InvestmentSnapshotTitledItem | dict[str, object],
    ) -> tuple[str, str]:
        """
        目的：兼容不同来源的亮点或风险条目对象。
        功能：统一抽取 `title` 和 `detail` 两个字段。
        实现逻辑：优先读取模型对象属性，若是字典则按键名读取。
        可调参数：`item` 为单个亮点或风险项。
        默认参数及原因：无法识别时返回空字符串，原因是上层会负责补位。
        """

        if isinstance(item, InvestmentSnapshotTitledItem):
            return str(item.title).strip(), str(item.detail).strip()
        if isinstance(item, dict):
            return str(item.get("title", "")).strip(), str(item.get("detail", "")).strip()
        return "", ""

    def _parse_financial_markdown(self, financial_source_markdown: str) -> ParsedFinancialTable:
        """
        目的：从 finance pack 的 Markdown 原文中解析快照页需要的财务基础数据。
        功能：定位“核心财务数据总表”，抽出最近 3 期列，并建立按指标名索引的值映射。
        实现逻辑：先找标题，再抓取其后的首张 Markdown 表，最后按“真实期间列 + 指标名称列”重建结构。
        可调参数：`financial_source_markdown` 为完整 finance pack 原文。
        默认参数及原因：只解析该章节下第一张核心表，原因是当前 v0.3 finance pack 已按此约定输出。
        """

        markdown = self._normalize_financial_markdown(financial_source_markdown)
        if not markdown:
            raise ValueError("financial_source_markdown must not be empty.")

        lines = markdown.splitlines()
        heading_index = next(
            (index for index, line in enumerate(lines) if self._is_core_financial_table_heading(line)),
            -1,
        )
        if heading_index < 0:
            raise ValueError(self._build_core_financial_heading_not_found_message(lines))

        table_lines: list[str] = []
        found_table = False
        for line in lines[heading_index + 1 :]:
            stripped = line.strip()
            if self._is_markdown_heading(stripped) and found_table:
                break
            if stripped.startswith("|"):
                table_lines.append(stripped)
                found_table = True
                continue
            if found_table and stripped:
                break

        if len(table_lines) < 3:
            raise ValueError("Failed to parse the markdown financial table from financial_source_markdown.")

        header_cells = self._parse_markdown_table_row(table_lines[0])
        period_indexes = [index for index, cell in enumerate(header_cells) if self._looks_like_period(cell)]
        if not period_indexes:
            raise ValueError("The core finance table does not contain any recognizable period columns.")

        label_column_index = max(0, period_indexes[0] - 1)
        raw_periods = [header_cells[index] for index in period_indexes]
        selected_indexes = list(range(max(0, len(raw_periods) - 3), len(raw_periods)))

        row_lookup: dict[str, list[str]] = {}
        for line in table_lines[2:]:
            if self._is_markdown_separator_row(line):
                continue
            cells = self._parse_markdown_table_row(line)
            if len(cells) <= label_column_index:
                continue
            label = cells[label_column_index].strip()
            if not label:
                continue
            row_lookup[label] = [
                self._normalize_text(cells[index] if index < len(cells) else "", fallback_text=FALLBACK_TEXT)
                for index in period_indexes
            ]

        if not row_lookup:
            raise ValueError("No usable data rows were parsed from the core finance table.")

        return ParsedFinancialTable(
            periods=raw_periods,
            display_periods=[raw_periods[index] for index in selected_indexes],
            selected_indexes=selected_indexes,
            unit_note=self._extract_financial_unit(markdown),
            row_lookup=row_lookup,
        )

    def _normalize_financial_markdown(self, financial_source_markdown: str | None) -> str:
        """
        目的：在进入财务解析前先做最小清洗。
        功能：去掉整体首尾空白，并兼容 agent 可能额外包上的成对 Markdown 代码块围栏。
        实现逻辑：只剥离最外层完整围栏，不改写正文内部结构，避免影响后续标题和表格解析。
        可调参数：`financial_source_markdown` 为工具接收到的财务原文字符串。
        默认参数及原因：只做最小清洗，原因是当前目标是提升兼容性而不是容忍任意自由格式输入。
        """

        markdown = str(financial_source_markdown or "").strip()
        if not markdown:
            return ""

        fenced_match = re.fullmatch(
            r"```(?:markdown|md)?\s*\n(?P<body>.*)\n```",
            markdown,
            flags=re.IGNORECASE | re.DOTALL,
        )
        if fenced_match:
            return fenced_match.group("body").strip()
        return markdown

    def _build_core_financial_heading_not_found_message(self, lines: list[str]) -> str:
        """
        目的：在缺少核心财务表标题时给出更可排查的错误信息。
        功能：把当前 Markdown 中能识别到的标题候选一并带出，帮助判断是 agent 改写了输入还是 pack 结构漂移。
        实现逻辑：先提取最多 6 条标题候选；若完全没有标题，再退回前几行非空文本作为上下文。
        可调参数：`lines` 为财务 Markdown 按行拆分后的结果。
        默认参数及原因：最多展示 6 条候选，原因是排查时需要足够上下文，但不应让错误信息过长。
        """

        heading_candidates = self._collect_heading_candidates(lines)
        if heading_candidates:
            return (
                "Failed to find the core finance table heading in financial_source_markdown. "
                f"Recognized heading candidates: {', '.join(heading_candidates)}"
            )

        preview_lines = [line.strip() for line in lines if line.strip()][:5]
        if preview_lines:
            return (
                "Failed to find the core finance table heading in financial_source_markdown. "
                f"Leading non-empty lines: {' | '.join(preview_lines)}"
            )
        return "Failed to find the core finance table heading in financial_source_markdown."

    def _collect_heading_candidates(self, lines: list[str]) -> list[str]:
        """
        目的：收集当前 Markdown 中可疑似标题的文本候选。
        功能：让解析失败时能快速看到 agent 实际传进来的章节结构。
        实现逻辑：遍历 Markdown 标题行，提取标题正文并按出现顺序保留前 6 条。
        可调参数：`lines` 为财务 Markdown 的逐行文本。
        默认参数及原因：只收前 6 条，原因是对单次排错已经足够且能避免错误消息冗长。
        """

        heading_candidates: list[str] = []
        for raw_line in lines:
            if not self._is_markdown_heading(raw_line):
                continue
            heading_text = self._extract_heading_text(raw_line)
            if not heading_text:
                continue
            heading_candidates.append(heading_text)
            if len(heading_candidates) >= 6:
                break
        return heading_candidates

    def _is_markdown_heading(self, line: str) -> bool:
        """
        目的：识别一行文本是否为 Markdown 标题。
        功能：统一给章节扫描和表格截断逻辑复用，避免只认 `## ` 这一种过窄格式。
        实现逻辑：允许 1 到 6 级标题，并兼容少量前置空白。
        可调参数：`line` 为单行 Markdown 文本。
        默认参数及原因：只匹配标准井号标题，原因是当前 finance pack 由 Markdown 模板稳定产出。
        """

        return bool(re.match(r"^\s{0,3}#{1,6}\s+", str(line or "")))

    def _extract_heading_text(self, line: str) -> str:
        """
        目的：把 Markdown 标题行规整为便于语义匹配的正文文本。
        功能：去掉标题井号、常见数字编号和多余空白，让“## 1. 核心财务数据总表”与“## 核心财务数据总表”统一处理。
        实现逻辑：先剥离标题前缀，再剥离常见中文和英文编号前缀，最后压缩空白。
        可调参数：`line` 为单行 Markdown 标题文本。
        默认参数及原因：仅移除常见编号前缀，原因是当前目标是覆盖真实 pack，而不是支持任意富文本标题格式。
        """

        normalized_line = " ".join(str(line or "").split()).strip()
        normalized_line = re.sub(r"^\s{0,3}#{1,6}\s*", "", normalized_line)
        normalized_line = re.sub(r"^[0-9一二三四五六七八九十]+[\.\-、，．\)]\s*", "", normalized_line)
        return normalized_line.strip()

    def _is_core_financial_table_heading(self, line: str) -> bool:
        """
        目的：识别 finance pack 中的“核心财务数据总表”标题。
        功能：兼容少量空白差异，避免因为标题格式轻微波动导致整张 PPT 失败。
        实现逻辑：统一压缩空白后，用包含式规则匹配关键标题词。
        可调参数：`line` 为单行 Markdown 文本。
        默认参数及原因：只放宽到当前仓库真实需要的范围，原因是过宽会误抓错表。
        """

        heading_text = self._extract_heading_text(line)
        return bool(heading_text) and "核心财务数据总表" in heading_text

    def _parse_markdown_table_row(self, line: str) -> list[str]:
        """
        目的：拆分单行 Markdown 表格文本。
        功能：把 `|` 分隔的标准表格行转成单元格列表。
        实现逻辑：去掉首尾管道符后按 `|` 切分，再对每格做去空白。
        可调参数：`line` 为单行表格文本。
        默认参数及原因：只处理标准 Markdown 表格，原因是 finance pack 由模板稳定产出。
        """

        return [cell.strip() for cell in line.strip().strip("|").split("|")]

    def _is_markdown_separator_row(self, line: str) -> bool:
        """
        目的：识别 Markdown 表格中的分隔行。
        功能：避免把 `| --- | --- |` 误当成数据行。
        实现逻辑：判断整行是否仅由 `|:- ` 这些字符构成。
        可调参数：`line` 为单行表格文本。
        默认参数及原因：用最小规则识别，原因是当前模板输出都是标准分隔行。
        """

        stripped = line.strip()
        return bool(stripped) and all(char in "|:- " for char in stripped)

    def _looks_like_period(self, value: str) -> bool:
        """
        目的：识别表头中的真实期间列。
        功能：兼容 `2025A`、`2025H1`、`2025Q3`、`TTM`、`FY-1` 等常见写法。
        实现逻辑：先做空白归一化，再用统一正则和年份补充规则匹配。
        可调参数：`value` 为表头单元格文本。
        默认参数及原因：只识别常见财务期间，原因是当前输出场景不需要更宽松的时间语义。
        """

        normalized = " ".join(str(value or "").split()).strip().lower()
        if not normalized:
            return False
        if self._PERIOD_PATTERN.fullmatch(normalized):
            return True
        return bool(re.fullmatch(r"(19|20)\d{2}", normalized))

    def _extract_financial_unit(self, markdown: str) -> str:
        """
        目的：从 finance pack 中提取币种和单位说明。
        功能：让快照页财务表头能尽量保留“人民币 / 百万元 / 千元”等真实口径。
        实现逻辑：扫描全文中包含“单位”“币种”“货币”的行，命中即清洗输出。
        可调参数：`markdown` 为 finance pack 原文。
        默认参数及原因：未命中时回落到“单位未披露”，原因是宁可明确缺失也不猜单位。
        """

        for raw_line in markdown.splitlines():
            line = " ".join(raw_line.split()).strip()
            if not line:
                continue
            if not any(keyword in line for keyword in ("单位", "币种", "货币")):
                continue
            cleaned_line = re.sub(r"^[>*#\\-\\s]+", "", line).strip()
            if cleaned_line:
                return cleaned_line
        return "单位未披露"

    def _build_snapshot_financial_rows(
        self,
        parsed_table: ParsedFinancialTable,
    ) -> list[InvestmentSnapshotFinancialRow]:
        """
        目的：生成快照页固定展示的 8 行财务指标。
        功能：优先取显式披露值，缺失时做有限推导，并在仍无法确定时填“缺乏信息”。
        实现逻辑：先取营业收入、净利润、总资产、股东权益等基础序列，再依序补各个目标行。
        可调参数：`parsed_table` 为已解析好的核心财务表结果。
        默认参数及原因：不做超出基础会计关系的复杂推导，原因是快照页更重视稳妥而不是覆盖率极限。
        """

        row_lookup = parsed_table.row_lookup
        revenue_series = self._pick_series_by_aliases(row_lookup, DIRECT_FINANCIAL_LABEL_ALIASES["营业收入"])
        net_profit_series = self._pick_series_by_aliases(row_lookup, SUPPORTING_FINANCIAL_LABEL_ALIASES["净利润"])
        total_assets_series = self._pick_series_by_aliases(row_lookup, DIRECT_FINANCIAL_LABEL_ALIASES["总资产"])
        equity_series = self._pick_series_by_aliases(row_lookup, SUPPORTING_FINANCIAL_LABEL_ALIASES["股东权益"])

        result_rows: list[InvestmentSnapshotFinancialRow] = []
        for label in REQUIRED_FINANCIAL_LABELS:
            display_values = self._select_display_values(
                self._pick_series_by_aliases(row_lookup, DIRECT_FINANCIAL_LABEL_ALIASES[label]),
                parsed_table.selected_indexes,
            )
            if not self._has_meaningful_values(display_values):
                if label == "营业收入增长率":
                    display_values = self._select_display_values(
                        self._compute_growth_series(revenue_series),
                        parsed_table.selected_indexes,
                    )
                elif label == "净利率":
                    display_values = self._select_display_values(
                        self._compute_ratio_series(net_profit_series, revenue_series),
                        parsed_table.selected_indexes,
                    )
                elif label == "净利润增长率":
                    display_values = self._select_display_values(
                        self._compute_growth_series(net_profit_series),
                        parsed_table.selected_indexes,
                    )
                elif label == "资产负债率":
                    display_values = self._select_display_values(
                        self._compute_debt_ratio_series(total_assets_series, equity_series),
                        parsed_table.selected_indexes,
                    )

            result_rows.append(
                InvestmentSnapshotFinancialRow(
                    label=label,
                    values=[
                        self._normalize_text(value, fallback_text=FALLBACK_TEXT)
                        for value in display_values
                    ],
                )
            )
        return result_rows

    def _pick_series_by_aliases(self, row_lookup: dict[str, list[str]], aliases: tuple[str, ...]) -> list[str]:
        """
        目的：按别名集合选取一整行财务序列。
        功能：兼容上游 finance pack 里同一指标的不同命名方式。
        实现逻辑：按别名优先级顺序逐个匹配，命中即返回对应整行。
        可调参数：`row_lookup` 为行映射，`aliases` 为当前指标允许的名称集合。
        默认参数及原因：未命中时返回空列表，原因是后续统一补缺和推导更简单。
        """

        for alias in aliases:
            if alias in row_lookup:
                return list(row_lookup[alias])
        return []

    def _select_display_values(self, series: list[str], selected_indexes: list[int]) -> list[str]:
        """
        目的：从全量期间序列中截取当前展示窗口。
        功能：把整行历史值压缩成快照页实际展示的最近 3 期。
        实现逻辑：按 `selected_indexes` 顺序取值，越界位置补“缺乏信息”。
        可调参数：`series` 为完整序列，`selected_indexes` 为展示列索引。
        默认参数及原因：统一补缺，原因是所有目标行都必须严格对齐表头期间。
        """

        if not series:
            return [FALLBACK_TEXT for _ in selected_indexes]
        return [series[index] if index < len(series) else FALLBACK_TEXT for index in selected_indexes]

    def _has_meaningful_values(self, values: list[str]) -> bool:
        """
        目的：判断一行展示值是否已包含有效内容。
        功能：区分“全为空占位”与“至少有一格是真实值”。
        实现逻辑：只要任意单元格不是 `缺乏信息` 就视为命中。
        可调参数：`values` 为当前展示值列表。
        默认参数及原因：显式文本即使不是数值也视为有效，原因是上游可能直接给出可读说明。
        """

        return any(value.strip() and value.strip() != FALLBACK_TEXT for value in values)

    def _compute_growth_series(self, series: list[str]) -> list[str]:
        """
        目的：在缺少显式增长率行时按相邻期间推导增长率。
        功能：基于上一期绝对值计算同比或期比变化率。
        实现逻辑：使用 `(本期 - 上期) / abs(上期)`，首期或分母缺失时填“缺乏信息”。
        可调参数：`series` 为需要计算增长率的绝对值序列。
        默认参数及原因：分母取绝对值，原因是净利润为负时也能保持方向表达稳定。
        """

        numeric_values = [self._parse_numeric_value(value) for value in series]
        growth_values: list[str] = []
        for index, current_value in enumerate(numeric_values):
            if index == 0:
                growth_values.append(FALLBACK_TEXT)
                continue
            previous_value = numeric_values[index - 1]
            if current_value is None or previous_value is None or previous_value == 0:
                growth_values.append(FALLBACK_TEXT)
                continue
            growth_values.append(self._format_percent_value((current_value - previous_value) / abs(previous_value) * 100))
        return growth_values

    def _compute_ratio_series(
        self,
        numerator_series: list[str],
        denominator_series: list[str],
    ) -> list[str]:
        """
        目的：在缺少显式比率时按分子分母序列推导百分比。
        功能：当前主要用于从净利润和营业收入推导净利率。
        实现逻辑：逐期解析数值并计算 `分子 / 分母 * 100`，分母无效时补缺。
        可调参数：`numerator_series` 和 `denominator_series` 为同期间对齐的序列。
        默认参数及原因：只做直接比率推导，原因是这类公式简单且不会引入额外口径歧义。
        """

        ratio_values: list[str] = []
        for index in range(max(len(numerator_series), len(denominator_series))):
            numerator = self._parse_numeric_value(numerator_series[index]) if index < len(numerator_series) else None
            denominator = self._parse_numeric_value(denominator_series[index]) if index < len(denominator_series) else None
            if numerator is None or denominator in (None, 0):
                ratio_values.append(FALLBACK_TEXT)
                continue
            ratio_values.append(self._format_percent_value(numerator / denominator * 100))
        return ratio_values

    def _compute_debt_ratio_series(
        self,
        total_assets_series: list[str],
        equity_series: list[str],
    ) -> list[str]:
        """
        目的：在缺少显式资产负债率时用总资产和股东权益做有限推导。
        功能：按基础会计恒等式估算 `(总资产 - 股东权益) / 总资产`。
        实现逻辑：逐期解析数值，分母无效时补缺。
        可调参数：`total_assets_series` 和 `equity_series` 为同期间序列。
        默认参数及原因：只使用最基础的平衡关系，原因是这样风险最低且解释清楚。
        """

        debt_ratio_values: list[str] = []
        for index in range(max(len(total_assets_series), len(equity_series))):
            total_assets = self._parse_numeric_value(total_assets_series[index]) if index < len(total_assets_series) else None
            equity = self._parse_numeric_value(equity_series[index]) if index < len(equity_series) else None
            if total_assets in (None, 0) or equity is None:
                debt_ratio_values.append(FALLBACK_TEXT)
                continue
            debt_ratio_values.append(self._format_percent_value((1 - equity / total_assets) * 100))
        return debt_ratio_values

    def _parse_numeric_value(self, value: str) -> float | None:
        """
        目的：把财务单元格文本转成可计算数值。
        功能：兼容千分位、括号负数、百分号和简单文字占位。
        实现逻辑：先过滤明显不可计算文本，再抽取首个数值并还原负号。
        可调参数：`value` 为单元格文本。
        默认参数及原因：解析失败返回 `None`，原因是避免误把无信息当作 0。
        """

        normalized = str(value or "").strip()
        if not normalized or normalized in {FALLBACK_TEXT, "-", "--", "---", "N/A", "n/a"}:
            return None
        if any(flag in normalized for flag in ("无信息", "缺乏信息", "未披露", "不适用", "无法")):
            return None

        is_negative = normalized.startswith("(") and normalized.endswith(")")
        numeric_source = normalized.strip("()").replace(",", "")
        match = re.search(r"[-+]?\d+(?:\.\d+)?", numeric_source)
        if match is None:
            return None
        parsed_value = float(match.group())
        if is_negative and parsed_value > 0:
            return -parsed_value
        return parsed_value

    def _format_percent_value(self, value: float) -> str:
        """
        目的：统一格式化工具推导出的百分比结果。
        功能：把浮点数转成一位小数的百分比文本。
        实现逻辑：四舍五入到一位小数后再拼接 `%`。
        可调参数：`value` 为待格式化的百分比数值。
        默认参数及原因：保留一位小数，原因是快照页更重视清晰可读而不是极高精度。
        """

        return f"{round(value, 1):.1f}%"

    def _build_presentation(
        self,
        *,
        company_name: str,
        slide_title: str,
        positioning_line: str,
        overview_summary: str,
        overview_product_items: list[InvestmentSnapshotOverviewProductItem],
        financial_table: ParsedFinancialTable,
        financial_rows: list[InvestmentSnapshotFinancialRow],
        highlight_items: list[InvestmentSnapshotTitledItem],
        risk_items: list[InvestmentSnapshotTitledItem],
    ) -> Presentation:
        """
        目的：构建完整的一页式 investment snapshot 演示文稿。
        功能：按固定 A4 横向布局绘制标题、概况、财务、亮点和风险五个区域。
        实现逻辑：创建 blank slide 后依次绘制各区域，所有风格和尺寸都在工具内统一控制。
        可调参数：公司名、标题、摘要、财务表和结构化条目。
        默认参数及原因：只生成单页 blank layout，原因是当前任务目标就是稳定导出一页快照。
        """

        presentation = Presentation()
        presentation.slide_width = SLIDE_WIDTH
        presentation.slide_height = SLIDE_HEIGHT

        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BACKGROUND

        self._add_header_area(slide, company_name, slide_title, positioning_line)
        self._add_overview_panel(slide, overview_summary, overview_product_items)
        self._add_financial_panel(slide, financial_table, financial_rows)
        self._add_highlight_panel(slide, highlight_items)
        self._add_risk_panel(slide, risk_items)
        return presentation

    def _add_header_area(
        self,
        slide,
        company_name: str,
        slide_title: str,
        positioning_line: str,
    ) -> None:
        """
        目的：绘制页面顶部标题和定位横幅。
        功能：展示“投资要点速览 - 公司名”主标题、金色分隔线和红色定位条。
        实现逻辑：用文本框和窄矩形组合成稳定页首层次。
        可调参数：公司名、页面主标题和定位短句。
        默认参数及原因：页首高度固定，原因是下方四个内容区都依赖这组锚点排版。
        """

        self._add_textbox(
            slide=slide,
            left=PAGE_LEFT,
            top=TITLE_TOP,
            width=PAGE_WIDTH,
            height=TITLE_HEIGHT,
            text=f"{slide_title} - {company_name}",
            font_size=20,
            bold=True,
            font_name=TITLE_FONT,
            color=TITLE_TEXT,
            align=PP_ALIGN.LEFT,
            margin_left=2,
            margin_right=2,
            margin_top=0,
            margin_bottom=0,
        )

        rule = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            PAGE_LEFT,
            TITLE_RULE_TOP,
            PAGE_WIDTH,
            Inches(0.03),
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = SECONDARY_GOLD
        rule.line.fill.background()

        self._add_textbox(
            slide=slide,
            left=PAGE_LEFT,
            top=BANNER_TOP,
            width=PAGE_WIDTH,
            height=BANNER_HEIGHT,
            text=positioning_line,
            font_size=13.5,
            bold=True,
            font_name=TITLE_FONT,
            color=RGBColor(255, 255, 255),
            fill_color=PRIMARY_RED,
            line_color=PRIMARY_RED,
            align=PP_ALIGN.LEFT,
            margin_left=10,
            margin_right=10,
            margin_top=0,
            margin_bottom=0,
        )

    def _add_overview_panel(
        self,
        slide,
        overview_summary: str,
        overview_product_items: list[InvestmentSnapshotOverviewProductItem],
    ) -> None:
        """
        目的：绘制左上角公司概况区。
        功能：展示一段摘要和 3 条产品项。
        实现逻辑：先画容器，再在文本框里按“摘要 + 产品项”顺序拼段落。
        可调参数：摘要文本和产品项列表。
        默认参数及原因：区域尺寸固定，原因是要与右侧财务区对齐。
        """

        left = PAGE_LEFT
        top = TOP_PANEL_TOP
        width = TOP_PANEL_WIDTH
        height = TOP_PANEL_HEIGHT

        self._add_panel_box(slide, left, top, width, height, "公司概况")
        textbox = self._create_text_container(
            slide,
            left + Inches(0.16),
            top + Inches(0.20),
            width - Inches(0.32),
            height - Inches(0.28),
        )
        text_frame = textbox.text_frame
        font_size = 10.8 if len(overview_summary) < 140 else 10.2
        self._add_plain_paragraph(text_frame, overview_summary, font_size, MUTED_TEXT, is_first=True, space_after=4)
        for item in overview_product_items:
            self._add_named_item_paragraph(text_frame, item, font_size, PRIMARY_RED, TITLE_TEXT, space_after=2)

    def _add_financial_panel(
        self,
        slide,
        financial_table: ParsedFinancialTable,
        financial_rows: list[InvestmentSnapshotFinancialRow],
    ) -> None:
        """
        目的：绘制右上角固定 8 行财务表。
        功能：展示最近 3 期期间列和固定指标行。
        实现逻辑：用一张 PowerPoint 表格承载数据，再单独给表头和边框做样式。
        可调参数：解析后的财务表元信息和 8 行指标结果。
        默认参数及原因：期间列数跟随解析结果，原因是工具需要兼容真实样本里 1 到 3 列的情况。
        """

        left = PAGE_LEFT + TOP_PANEL_WIDTH + PANEL_GAP
        top = TOP_PANEL_TOP
        width = TOP_PANEL_WIDTH
        height = TOP_PANEL_HEIGHT
        period_count = len(financial_table.display_periods)

        self._add_panel_box(slide, left, top, width, height, "财务数据")

        table_shape = slide.shapes.add_table(
            rows=1 + len(financial_rows),
            cols=1 + period_count,
            left=left + Inches(0.12),
            top=top + Inches(0.20),
            width=width - Inches(0.24),
            height=height - Inches(0.30),
        )
        table = table_shape.table
        first_column_width = Inches(1.72)
        other_column_width = int((width - Inches(0.24) - first_column_width) / max(1, period_count))
        table.columns[0].width = first_column_width
        for column_index in range(1, 1 + period_count):
            table.columns[column_index].width = other_column_width

        header_labels = [f"项目（{financial_table.unit_note}）"] + financial_table.display_periods
        for column_index, header_label in enumerate(header_labels):
            self._set_table_cell_text(
                table.cell(0, column_index),
                header_label,
                font_size=9.5,
                bold=True,
                align=PP_ALIGN.CENTER,
                fill_color=LABEL_BG,
            )

        for row_index, row_payload in enumerate(financial_rows, start=1):
            fill_color = ROW_ODD_BG if row_index % 2 == 1 else ROW_EVEN_BG
            self._set_table_cell_text(
                table.cell(row_index, 0),
                row_payload.label,
                font_size=9.8,
                bold=True,
                align=PP_ALIGN.LEFT,
                fill_color=fill_color,
            )
            for value_index, value in enumerate(row_payload.values, start=1):
                self._set_table_cell_text(
                    table.cell(row_index, value_index),
                    value,
                    font_size=9.6,
                    bold=False,
                    align=PP_ALIGN.RIGHT,
                    fill_color=fill_color,
                )

    def _add_highlight_panel(
        self,
        slide,
        highlight_items: list[InvestmentSnapshotTitledItem],
    ) -> None:
        """
        目的：绘制页面中下部的投资亮点区。
        功能：承载 3 到 5 条亮点条目。
        实现逻辑：先画容器，再用富文本段落按“标题加粗 + 说明”方式输出。
        可调参数：亮点条目列表。
        默认参数及原因：区域横向铺满页面，原因是亮点往往条目较多且需要更宽阅读面。
        """

        left = PAGE_LEFT
        top = HIGHLIGHT_TOP
        width = PAGE_WIDTH
        height = HIGHLIGHT_HEIGHT

        self._add_panel_box(slide, left, top, width, height, "投资亮点")
        textbox = self._create_text_container(
            slide,
            left + Inches(0.16),
            top + Inches(0.18),
            width - Inches(0.32),
            height - Inches(0.26),
        )
        text_frame = textbox.text_frame
        font_size = self._fit_titled_items_font_size(highlight_items, base_size=10.6, min_size=9.6, chars_per_line=105)
        for index, item in enumerate(highlight_items):
            self._add_titled_item_paragraph(
                text_frame,
                item,
                font_size=font_size,
                title_color=PRIMARY_RED,
                detail_color=TITLE_TEXT,
                is_first=index == 0,
                space_after=2,
            )

    def _add_risk_panel(
        self,
        slide,
        risk_items: list[InvestmentSnapshotTitledItem],
    ) -> None:
        """
        目的：绘制页面底部的投资风险区。
        功能：承载 1 到 2 条风险条目，并保持可快速扫读。
        实现逻辑：复用统一的标题段落样式，但采用更紧凑的区域高度和字号。
        可调参数：风险条目列表。
        默认参数及原因：底部区域高度固定偏小，原因是当前风险项条数上限较低。
        """

        left = PAGE_LEFT
        top = RISK_TOP
        width = PAGE_WIDTH
        height = RISK_HEIGHT

        self._add_panel_box(slide, left, top, width, height, "投资风险")
        textbox = self._create_text_container(
            slide,
            left + Inches(0.16),
            top + Inches(0.14),
            width - Inches(0.32),
            height - Inches(0.20),
        )
        text_frame = textbox.text_frame
        font_size = self._fit_titled_items_font_size(risk_items, base_size=10.4, min_size=9.4, chars_per_line=118)
        for index, item in enumerate(risk_items):
            self._add_titled_item_paragraph(
                text_frame,
                item,
                font_size=font_size,
                title_color=PRIMARY_RED,
                detail_color=TITLE_TEXT,
                is_first=index == 0,
                space_after=1,
            )

    def _fit_titled_items_font_size(
        self,
        items: list[InvestmentSnapshotTitledItem],
        *,
        base_size: float,
        min_size: float,
        chars_per_line: float,
    ) -> float:
        """
        目的：根据条目文字量粗略收缩亮点或风险区字号。
        功能：避免长条目在固定高度里明显溢出。
        实现逻辑：按加权字符数估算行数，内容越多字号越小，但不会低于最小值。
        可调参数：条目列表、基准字号、最小字号和每行字符容量估计。
        默认参数及原因：只做启发式缩放，原因是 PowerPoint 不提供可靠的写后测高回读。
        """

        total_chars = sum(len(item.title) * 1.4 + len(item.detail) for item in items)
        if total_chars <= chars_per_line * 3.2:
            return base_size
        shrink_steps = math.ceil((total_chars - chars_per_line * 3.2) / max(1.0, chars_per_line * 0.9))
        return max(min_size, round(base_size - 0.2 * shrink_steps, 1))

    def _add_panel_box(self, slide, left, top, width, height, label: str) -> None:
        """
        目的：绘制通用内容面板容器。
        功能：输出带浅底色、金色边框和顶部标签的面板。
        实现逻辑：先画背景框，再叠加小标签文本框。
        可调参数：位置尺寸和标签文案。
        默认参数及原因：统一面板风格，原因是快照页需要整体视觉一致。
        """

        panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
        panel.fill.solid()
        panel.fill.fore_color.rgb = PANEL_BG
        panel.line.color.rgb = PANEL_BORDER
        panel.line.width = Pt(1.2)

        self._add_textbox(
            slide=slide,
            left=left + Inches(0.14),
            top=top - Inches(0.14),
            width=Inches(1.60),
            height=Inches(0.28),
            text=label,
            font_size=12.2,
            bold=True,
            font_name=TITLE_FONT,
            color=PRIMARY_RED,
            fill_color=LABEL_BG,
            line_color=LABEL_BG,
            align=PP_ALIGN.CENTER,
            margin_left=0,
            margin_right=0,
            margin_top=0,
            margin_bottom=0,
        )

    def _create_text_container(self, slide, left, top, width, height):
        """
        目的：创建一个纯文本承载容器。
        功能：统一设置自动换行、内边距和垂直对齐方式。
        实现逻辑：使用透明文本框作为容器，再返回给后续段落写入方法复用。
        可调参数：位置尺寸。
        默认参数及原因：默认启用自动换行并去掉边框，原因是正文区只需要内容不需要额外描边。
        """

        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.fill.background()
        textbox.line.fill.background()
        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.margin_left = Pt(0)
        text_frame.margin_right = Pt(0)
        text_frame.margin_top = Pt(0)
        text_frame.margin_bottom = Pt(0)
        return textbox

    def _add_plain_paragraph(
        self,
        text_frame,
        text: str,
        font_size: float,
        color: RGBColor,
        *,
        is_first: bool,
        space_after: int,
    ) -> None:
        """
        目的：向文本容器写入普通正文段落。
        功能：统一控制摘要类文本的字号、颜色和段后距。
        实现逻辑：首段复用现有段落，其余情况新建段落并写入单个 run。
        可调参数：文本框、正文文本、字号、颜色和段后距。
        默认参数及原因：首段复用是为了避免空段残留，原因是 pptx 文本框默认自带一个空段。
        """

        paragraph = text_frame.paragraphs[0] if is_first else text_frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(space_after)
        run = paragraph.add_run()
        run.text = text
        self._apply_run_font(run, font_size, False, PRIMARY_FONT, color)

    def _add_named_item_paragraph(
        self,
        text_frame,
        item: InvestmentSnapshotOverviewProductItem,
        font_size: float,
        title_color: RGBColor,
        detail_color: RGBColor,
        *,
        space_after: int,
    ) -> None:
        """
        目的：向概况区写入“产品名加粗 + 描述”段落。
        功能：让产品条目具备清晰的标题层和说明层。
        实现逻辑：同一段中先写标题 run，再写说明 run。
        可调参数：文本框、产品项、字号、标题色、说明色和段后距。
        默认参数及原因：标题和说明共段落，原因是这样最节省纵向空间。
        """

        paragraph = text_frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(space_after)
        title_run = paragraph.add_run()
        title_run.text = f"{item.name}："
        self._apply_run_font(title_run, font_size, True, PRIMARY_FONT, title_color)
        detail_run = paragraph.add_run()
        detail_run.text = item.description
        self._apply_run_font(detail_run, font_size, False, PRIMARY_FONT, detail_color)

    def _add_titled_item_paragraph(
        self,
        text_frame,
        item: InvestmentSnapshotTitledItem,
        font_size: float,
        title_color: RGBColor,
        detail_color: RGBColor,
        *,
        is_first: bool,
        space_after: int,
    ) -> None:
        """
        目的：写入“标题加粗 + 说明”的亮点或风险段落。
        功能：保持亮点和风险区统一的可扫描结构。
        实现逻辑：与概况产品项类似，但标题改用更短的分析标签。
        可调参数：文本框、条目、字号、颜色、首段标记和段后距。
        默认参数及原因：标题与详情放在同段，原因是当前单页高度有限。
        """

        paragraph = text_frame.paragraphs[0] if is_first else text_frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(space_after)
        title_run = paragraph.add_run()
        title_run.text = f"{item.title}："
        self._apply_run_font(title_run, font_size, True, PRIMARY_FONT, title_color)
        detail_run = paragraph.add_run()
        detail_run.text = item.detail
        self._apply_run_font(detail_run, font_size, False, PRIMARY_FONT, detail_color)

    def _add_textbox(
        self,
        *,
        slide,
        left,
        top,
        width,
        height,
        text: str,
        font_size: float,
        bold: bool,
        font_name: str = PRIMARY_FONT,
        color: RGBColor = TITLE_TEXT,
        fill_color: RGBColor | None = None,
        line_color: RGBColor | None = None,
        align: PP_ALIGN = PP_ALIGN.LEFT,
        margin_left: int = 4,
        margin_right: int = 4,
        margin_top: int = 2,
        margin_bottom: int = 2,
    ):
        """
        目的：封装常用文本框创建逻辑。
        功能：统一处理填充、边框、内边距和首段文本样式。
        实现逻辑：先建文本框，再按参数配置样式并写入一段文本。
        可调参数：位置尺寸、文本、字体样式、填充色、边框色和对齐方式。
        默认参数及原因：边距取较小值，原因是当前快照页空间紧凑。
        """

        textbox = slide.shapes.add_textbox(left, top, width, height)
        if fill_color is None:
            textbox.fill.background()
        else:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = fill_color
        if line_color is None:
            textbox.line.fill.background()
        else:
            textbox.line.color.rgb = line_color
            textbox.line.width = Pt(0.8)

        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.margin_left = Pt(margin_left)
        text_frame.margin_right = Pt(margin_right)
        text_frame.margin_top = Pt(margin_top)
        text_frame.margin_bottom = Pt(margin_bottom)

        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        self._apply_run_font(run, font_size, bold, font_name, color)
        return textbox

    def _set_table_cell_text(
        self,
        cell,
        text: str,
        *,
        font_size: float,
        bold: bool,
        align: PP_ALIGN,
        fill_color: RGBColor,
    ) -> None:
        """
        目的：统一设置表格单元格文本和样式。
        功能：控制字体、对齐、底色和边距，减少 PowerPoint 默认样式干扰。
        实现逻辑：清空原文本后写入一个 run，再设置 fill 和 margin。
        可调参数：单元格对象、文本、字号、粗细、对齐和底色。
        默认参数及原因：数值列右对齐，原因是财务表阅读时更容易横向比较。
        """

        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if align == PP_ALIGN.RIGHT:
            cell.margin_left = Pt(1.5)
            cell.margin_right = Pt(3.2)
        elif align == PP_ALIGN.LEFT:
            cell.margin_left = Pt(3.2)
            cell.margin_right = Pt(1.5)
        else:
            cell.margin_left = Pt(2.0)
            cell.margin_right = Pt(2.0)
        cell.margin_top = Pt(1.2)
        cell.margin_bottom = Pt(1.2)

        text_frame = cell.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        self._apply_run_font(run, font_size, bold, PRIMARY_FONT, TITLE_TEXT)

    def _apply_run_font(
        self,
        run,
        font_size: float,
        bold: bool,
        font_name: str,
        color: RGBColor,
    ) -> None:
        """
        目的：统一封装 run 级字体样式设置。
        功能：保证整页文本的字体、字色和粗细风格一致。
        实现逻辑：所有文本写入最终都经过这个方法落样式。
        可调参数：run 对象、字号、粗细、字体名和颜色。
        默认参数及原因：样式由调用方显式传入，原因是各分区层级不同但实现方式相同。
        """

        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
